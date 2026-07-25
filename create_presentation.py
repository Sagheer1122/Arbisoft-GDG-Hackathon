import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_professional_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    PURPLE_PRIMARY = RGBColor(81, 66, 197)   # #5142C5
    PURPLE_DARK = RGBColor(22, 22, 42)      # #16162A
    PURPLE_LIGHT = RGBColor(237, 233, 254)  # #EDE9FE
    WHITE = RGBColor(255, 255, 255)
    GRAY_TEXT = RGBColor(112, 112, 128)     # #707080
    GRAY_LIGHT = RGBColor(247, 247, 251)    # #F7F7FB
    BORDER_COLOR = RGBColor(231, 231, 240)  # #E7E7F0
    GOLD = RGBColor(245, 158, 11)           # #F59E0B
    GREEN = RGBColor(5, 150, 105)          # #059669

    blank_layout = prs.slide_layouts[6]

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header(slide, title_text, category_text="NURSEFLOW ENTERPRISE PLATFORM"):
        # Top Accent Line
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.04))
        accent.fill.solid()
        accent.fill.fore_color.rgb = PURPLE_PRIMARY
        accent.line.fill.background()

        # Category Pill
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = PURPLE_PRIMARY

        # Main Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.733), Inches(0.7))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = PURPLE_DARK

    def add_footer(slide, current_slide, total_slides=12):
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.4))
        tf = footer_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"NurseFlow Platform Presentation  |  Slide {current_slide} of {total_slides}  |  Arbisoft-GDG Hackathon 2026"
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY_TEXT

    # ==========================================
    # SLIDE 1: Title Slide (Hero Dark Theme)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, PURPLE_DARK)

    # Decorative Side Accent Bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PURPLE_PRIMARY
    bar.line.fill.background()

    # Glass Card Container
    hero_card = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(11.0), Inches(5.0))
    hero_card.fill.solid()
    hero_card.fill.fore_color.rgb = RGBColor(30, 30, 55)
    hero_card.line.color.rgb = PURPLE_PRIMARY
    hero_card.line.width = Pt(2)

    tf1 = hero_card.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.5)
    tf1.margin_top = Inches(0.5)

    p0 = tf1.paragraphs[0]
    p0.text = "🏥 NURSEFLOW"
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = WHITE

    p1 = tf1.add_paragraph()
    p1.text = "Smart Roster Management & Gemma 4 AI Clinical Training Platform"
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = GOLD
    p1.space_before = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "An enterprise full-stack healthcare solution combining automated duty scheduling, AI patient roleplay simulation, and nurse shift-break wellness lounge mini-games."
    p2.font.size = Pt(14)
    p2.font.color.rgb = PURPLE_LIGHT
    p2.space_before = Pt(16)

    # Badges Row inside Hero
    badges_box = slide1.shapes.add_textbox(Inches(1.7), Inches(4.8), Inches(10.0), Inches(0.8))
    tf_b = badges_box.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = "✦ React 18 & Vite   ✦ Node.js & Express   ✦ Prisma ORM & SQLite   ✦ Gemma 4 AI Engine   ✦ Socket.IO Push"
    p_b.font.size = Pt(12)
    p_b.font.bold = True
    p_b.font.color.rgb = WHITE

    # Presenter Footer
    fbox = slide1.shapes.add_textbox(Inches(1.2), Inches(6.5), Inches(11.0), Inches(0.5))
    ftf = fbox.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fp.text = "GitHub: Sagheer1122/Arbisoft-GDG-Hackathon  |  Live App: http://localhost:5173"
    fp.font.size = Pt(11)
    fp.font.color.rgb = GRAY_TEXT

    # ==========================================
    # SLIDE 2: Executive Vision & Challenges
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "Executive Overview: Solving Healthcare Workforce Fatigue")
    add_footer(slide2, 2)

    # Card 1: Hospital Challenges
    card1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    card1.fill.solid()
    card1.fill.fore_color.rgb = PURPLE_LIGHT
    card1.line.color.rgb = PURPLE_PRIMARY
    tf1 = card1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.3)
    tf1.margin_top = Inches(0.3)

    p = tf1.paragraphs[0]
    p.text = "🚨 Critical Healthcare Challenges"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PURPLE_DARK

    challenges = [
        "Scheduling Friction: Shift collisions, manual swap delays, and text overlap on duty rosters.",
        "Clinical Nurse Burnout: Exhaustion during intense 12-hour duty shifts with insufficient mental reset.",
        "Communication Practice Gap: Lack of realistic, safe scenarios for practicing high-stress patient interactions.",
        "Communication Delays: Slow emergency alert broadcasts across hospital wards."
    ]
    for ch in challenges:
        p_sub = tf1.add_paragraph()
        p_sub.text = "• " + ch
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = PURPLE_DARK
        p_sub.space_before = Pt(10)

    # Card 2: NurseFlow Solution
    card2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0))
    card2.fill.solid()
    card2.fill.fore_color.rgb = PURPLE_PRIMARY
    card2.line.color.rgb = PURPLE_PRIMARY
    tf2 = card2.text_frame
    tf2.word_wrap = True
    tf2.margin_left = Inches(0.3)
    tf2.margin_top = Inches(0.3)

    p_sol = tf2.paragraphs[0]
    p_sol.text = "🏥 The NurseFlow Solution"
    p_sol.font.size = Pt(18)
    p_sol.font.bold = True
    p_sol.font.color.rgb = WHITE

    solutions = [
        "Smart Roster Intelligence: Color-coded shifts (Morning, Evening, Night) with zero text overlap.",
        "Gemma 4 AI Patient Simulator: Educational roleplay scenarios with dynamic emotion tracking.",
        "Shift Break Lounge Games: 100% non-medical relaxation games (Tick & Cross) for 10-min breaks.",
        "Socket.IO Push Dispatch: Instant emergency ward broadcasts and 1-click swap approvals."
    ]
    for sol in solutions:
        p_sub2 = tf2.add_paragraph()
        p_sub2.text = "✓ " + sol
        p_sub2.font.size = Pt(12)
        p_sub2.font.color.rgb = WHITE
        p_sub2.space_before = Pt(10)

    # ==========================================
    # SLIDE 3: System Architecture & Ecosystem
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "System Architecture & Full-Stack Technology Stack")
    add_footer(slide3, 3)

    layers = [
        ("Frontend Single-Page App", "React 18 + Vite, TypeScript, Tailwind CSS, Lucide Icons, Glassmorphic Design System", PURPLE_PRIMARY),
        ("State Management & Router", "React Context API (AuthContext, SocketContext), React Router v6 Guarded Routes", PURPLE_DARK),
        ("Backend REST API", "Node.js + Express Server, JWT Bearer Token Security, Bcrypt Password Hashing Middleware", GREEN),
        ("Database & Persistence", "Prisma ORM with SQLite (dev.db), Seeded Scenarios, Roster Tables, and User Schemas", PURPLE_PRIMARY),
        ("Gemma 4 AI & Real-Time Engine", "Google Gemma 4 AI Roleplay Engine, 8-Score Analysis, Minimax Game AI, Socket.IO WebSockets", GOLD)
    ]

    top_pos = 1.6
    for title, desc, border_col in layers:
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.733), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = PURPLE_LIGHT
        card.line.color.rgb = border_col
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.12)

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(15)
        p1.font.bold = True
        p1.font.color.rgb = PURPLE_DARK

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = GRAY_TEXT

        top_pos += 1.02

    # ==========================================
    # SLIDE 4: Core Module 1 — Roster Calendar
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "Core Module 1: Smart Duty Roster Calendar Engine")
    add_footer(slide4, 4)

    # 4 Shift Cards Grid
    shift_types = [
        ("🌅 Morning Shift", "7:00 AM – 3:00 PM", "Soft Emerald Tint", "Assigned to Ward Triage & Morning Vitals", GREEN),
        ("🌇 Evening Shift", "3:00 PM – 11:00 PM", "Soft Amber Tint", "Ward Medication Rounds & Intake", GOLD),
        ("🌙 Night Shift", "11:00 PM – 7:00 AM", "Soft Purple Tint", "Overnight Monitoring & Vital Checks every 2 hrs", PURPLE_PRIMARY),
        ("🏖️ Off Duty", "No Shift Scheduled", "Soft Slate Tint", "Rest & Scheduled Downtime between shifts", GRAY_TEXT)
    ]

    left_pos = 0.8
    for title, hours, style, desc, color in shift_types:
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.6), Inches(2.75), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = color
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color

        p_hrs = tf.add_paragraph()
        p_hrs.text = hours
        p_hrs.font.size = Pt(12)
        p_hrs.font.bold = True
        p_hrs.font.color.rgb = PURPLE_DARK
        p_hrs.space_before = Pt(8)

        p_sty = tf.add_paragraph()
        p_sty.text = "Style: " + style
        p_sty.font.size = Pt(10)
        p_sty.font.color.rgb = GRAY_TEXT
        p_sty.space_before = Pt(6)

        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(11)
        p_desc.font.color.rgb = PURPLE_DARK
        p_desc.space_before = Pt(12)

        left_pos += 2.98

    # ==========================================
    # SLIDE 5: Core Module 2 — Gemma 4 AI Patient Simulator
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "Core Module 2: Gemma 4 AI Patient Communication Simulator")
    add_footer(slide5, 5)

    col1 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    col1.fill.solid()
    col1.fill.fore_color.rgb = PURPLE_DARK
    tf5_1 = col1.text_frame
    tf5_1.word_wrap = True
    tf5_1.margin_left = Inches(0.3)
    tf5_1.margin_top = Inches(0.3)

    p = tf5_1.paragraphs[0]
    p.text = "🤖 AI Persona Roleplay Engine"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GOLD

    sim_features = [
        "Pre-Seeded Scenarios: Anxious Patient, Angry Family Member, Confused Elderly Patient, Non-Cooperative Patient.",
        "Dynamic Character Emotion Shift: Evaluates nurse empathy signals ('understand', 'listen', 'help') to shift tone.",
        "Emotion Shift Pipeline: Hostile ➔ Frustrated ➔ Worried ➔ Calm ➔ Reassured.",
        "Educational Tooling: Safe roleplay environment to practice high-stress conversations."
    ]
    for sf in sim_features:
        p_sub = tf5_1.add_paragraph()
        p_sub.text = "• " + sf
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = WHITE
        p_sub.space_before = Pt(10)

    col2 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(5.6), Inches(5.0))
    col2.fill.solid()
    col2.fill.fore_color.rgb = PURPLE_LIGHT
    col2.line.color.rgb = PURPLE_PRIMARY
    tf5_2 = col2.text_frame
    tf5_2.word_wrap = True
    tf5_2.margin_left = Inches(0.3)
    tf5_2.margin_top = Inches(0.3)

    p_eval = tf5_2.paragraphs[0]
    p_eval.text = "📊 8-Competency Evaluation Engine"
    p_eval.font.size = Pt(18)
    p_eval.font.bold = True
    p_eval.font.color.rgb = PURPLE_DARK

    eval_scores = [
        "1. Empathy & Compassion Score (0-100)",
        "2. Active Listening & Response Depth",
        "3. Communication Clarity & Terminology",
        "4. Professionalism & Boundary Setting",
        "5. Emotional Intelligence & De-escalation",
        "6. Patient Engagement & Reassurance",
        "7. Confidence & Clinical Support",
        "8. Detailed Quote Highlights & Actionable Advice"
    ]
    for es in eval_scores:
        p_es = tf5_2.add_paragraph()
        p_es.text = "✓ " + es
        p_es.font.size = Pt(11)
        p_es.font.color.rgb = PURPLE_DARK
        p_es.space_before = Pt(6)

    # ==========================================
    # SLIDE 6: 8-Competency Scoring Radar
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "8-Competency AI Clinical Evaluation Metrics")
    add_footer(slide6, 6)

    metrics = [
        ("Empathy & Compassion", "Validating patient feelings & active care", GREEN),
        ("Active Listening", "Addressing patient questions without skipping details", PURPLE_PRIMARY),
        ("Communication Clarity", "Explaining medical terms in simple language", GOLD),
        ("Professionalism", "Maintaining calm composure under pressure", PURPLE_DARK),
        ("De-escalation Capability", "Neutralizing angry or confused situations", GREEN),
        ("Patient Engagement", "Involving patient in care decision-making", PURPLE_PRIMARY),
        ("Emotional Intelligence", "Recognizing emotional cues & tone shifts", GOLD),
        ("Clinical Confidence", "Reassuring patient with clear healthcare guidance", PURPLE_DARK)
    ]

    left = 0.8
    top = 1.6
    for i, (m_title, m_desc, col) in enumerate(metrics):
        box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.6), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = col
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.12)

        p = tf.paragraphs[0]
        p.text = f"{i+1}. {m_title}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = col

        p_d = tf.add_paragraph()
        p_d.text = m_desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = GRAY_TEXT
        p_d.space_before = Pt(4)

        top += 1.25
        if i == 3:
            left = 6.8
            top = 1.6

    # ==========================================
    # SLIDE 7: Core Module 3 — AI Break Game
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "Core Module 3: AI Tick & Cross (Tic-Tac-Toe ❌⭕) Break Game")
    add_footer(slide7, 7)

    box7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
    box7.fill.solid()
    box7.fill.fore_color.rgb = WHITE
    box7.line.color.rgb = PURPLE_PRIMARY
    box7.line.width = Pt(2)

    tf7 = box7.text_frame
    tf7.word_wrap = True
    tf7.margin_left = Inches(0.3)
    tf7.margin_top = Inches(0.3)

    p7 = tf7.paragraphs[0]
    p7.text = "🎮 Nurse Shift Break Relaxation Lounge"
    p7.font.size = Pt(20)
    p7.font.bold = True
    p7.font.color.rgb = PURPLE_PRIMARY

    game_bullets = [
        "100% Non-Medical Relaxation: Pure fun mini-game designed for 10-minute shift breaks to help nurses feel refreshed.",
        "Gemma 4 AI Opponent: Nurse plays ❌ or ⭕ against Gemma 4 AI in real-time.",
        "3 Difficulty Modes: Easy (Casual Chill), Medium (Smart Rival), and Unbeatable (Minimax Algorithm).",
        "Minimax Decision Tree Algorithm: Evaluates move trees up to depth 6 for optimal move selection.",
        "Scoreboard & Wellness Points: Tracks Nurse Wins, AI Wins, Draws, and awards Wellness Points to NurseGameScore model.",
        "Live 10-Minute Break Timer: Built-in countdown timer ensures shift breaks stay structured."
    ]
    for gb in game_bullets:
        p_gb = tf7.add_paragraph()
        p_gb.text = "• " + gb
        p_gb.font.size = Pt(13)
        p_gb.font.color.rgb = PURPLE_DARK
        p_gb.space_before = Pt(10)

    # ==========================================
    # SLIDE 8: Shift Swaps, Leave & Socket Alerts
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "Core Module 4: Shift Swaps, Leave Portal & Socket.IO Push Alerts")
    add_footer(slide8, 8)

    cards = [
        ("🔄 Shift Swap Requests", "Nurse-to-nurse duty exchange requests with peer selection, reason tracking, and 1-click admin approval workflows.", PURPLE_PRIMARY),
        ("🏖️ Leave Request Portal", "Annual, Sick, and Emergency leave applications with date range selection and status badges (Pending, Approved, Rejected).", GOLD),
        ("⚡ Socket.IO Push Alerts", "Bi-directional WebSocket event engine for instant emergency ward broadcasts and unread notification badge updates.", GREEN)
    ]

    left_pos = 0.8
    for title, desc, col in cards:
        box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_pos), Inches(1.6), Inches(3.64), Inches(5.0))
        box.fill.solid()
        box.fill.fore_color.rgb = PURPLE_LIGHT
        box.line.color.rgb = col
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = PURPLE_DARK

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = PURPLE_DARK
        p_d.space_before = Pt(14)

        left_pos += 4.04

    # ==========================================
    # SLIDE 9: Relational Database Architecture
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "Relational Database Architecture (Prisma ORM & SQLite)")
    add_footer(slide9, 9)

    models = [
        ("User", "id, name, email, passwordHash, role, employeeId, departmentId", PURPLE_PRIMARY),
        ("Roster", "id, nurseId, shiftId, departmentId, date, status, notes", PURPLE_DARK),
        ("CommunicationScenario", "id, title, category, description, characterRole, difficulty", GREEN),
        ("CommunicationSession", "id, userId, scenarioId, status, overallScore", GOLD),
        ("CommunicationAnalysis", "id, sessionId, overallScore, empathyScore, strengths, feedback", PURPLE_PRIMARY),
        ("NurseGameScore", "id, userId, gameType, nurseWins, aiWins, draws, pointsEarned", GREEN)
    ]

    left = 0.8
    top = 1.6
    for i, (m_name, m_fields, col) in enumerate(models):
        card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.6), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = col
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = f"Model: {m_name}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = col

        p_f = tf.add_paragraph()
        p_f.text = f"Fields: {m_fields}"
        p_f.font.size = Pt(11)
        p_f.font.color.rgb = GRAY_TEXT
        p_f.space_before = Pt(6)

        top += 1.65
        if i == 2:
            left = 6.8
            top = 1.6

    # ==========================================
    # SLIDE 10: Complete REST API Infrastructure
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "Complete Backend REST API Infrastructure")
    add_footer(slide10, 10)

    apis = [
        ("Authentication (/api/auth)", "POST /login, POST /register, GET /me"),
        ("Staff Management (/api/users)", "GET /, GET /:id, PATCH /:id"),
        ("Duty Rosters (/api/rosters)", "GET /, POST /, PATCH /:id, DELETE /:id"),
        ("Leave & Swaps (/api/leave-requests & /api/shift-swaps)", "GET /, POST /, PATCH /:id/approve, PATCH /:id/reject"),
        ("AI Patient Simulator (/api/communication-simulator)", "GET /scenarios, POST /sessions, POST /messages, POST /end"),
        ("AI Break Games (/api/games)", "POST /tic-tac-toe/move, POST /tic-tac-toe/score, GET /score")
    ]

    top_pos = 1.6
    for title, routes in apis:
        card = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top_pos), Inches(11.733), Inches(0.75))
        card.fill.solid()
        card.fill.fore_color.rgb = PURPLE_LIGHT
        card.line.color.rgb = PURPLE_PRIMARY

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.12)

        p = tf.paragraphs[0]
        p.text = title + "  ➔  " + routes
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PURPLE_DARK

        top_pos += 0.85

    # ==========================================
    # SLIDE 11: Design System & Avatars
    # ==========================================
    slide11 = prs.slides.add_slide(blank_layout)
    add_header(slide11, "Glassmorphic Design System & Photorealistic Visuals")
    add_footer(slide11, 11)

    box11 = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.0))
    box11.fill.solid()
    box11.fill.fore_color.rgb = PURPLE_DARK
    tf11 = box11.text_frame
    tf11.word_wrap = True
    tf11.margin_left = Inches(0.3)
    tf11.margin_top = Inches(0.3)

    p11 = tf11.paragraphs[0]
    p11.text = "🎨 Premium Design Tokens & Assets"
    p11.font.size = Pt(20)
    p11.font.bold = True
    p11.font.color.rgb = GOLD

    design_points = [
        "Curated HSL Color Tokens: Deep Purple (#5142C5), Dark Navy Mesh (#16162A), Soft Light Purple Accent (#EDE9FE).",
        "Glassmorphic Aesthetics: Liquid glass buttons (bg-white/15 backdrop-blur-md) and subtle scale hover micro-animations.",
        "Photorealistic 8K Nurse Avatars: Generated realistic portraits of BOTH a female nurse wearing full hijab and a male nurse together side-by-side (no cartoon graphics).",
        "Responsive Navigation: Responsive Sidebar, Topbar, and Mobile Bottom Nav bar across all viewports."
    ]
    for dp in design_points:
        p_dp = tf11.add_paragraph()
        p_dp.text = "✓ " + dp
        p_dp.font.size = Pt(13)
        p_dp.font.color.rgb = WHITE
        p_dp.space_before = Pt(12)

    # ==========================================
    # SLIDE 12: Conclusion & Project Links
    # ==========================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12, PURPLE_DARK)

    tbox12 = slide12.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.5))
    tf12 = tbox12.text_frame
    tf12.word_wrap = True

    p = tf12.paragraphs[0]
    p.text = "Thank You! Q&A & Demonstration"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    p_sub = tf12.add_paragraph()
    p_sub.text = "NurseFlow — Empowering Healthcare Staff Through AI Roster Intelligence"
    p_sub.font.size = Pt(18)
    p_sub.font.bold = True
    p_sub.font.color.rgb = GOLD
    p_sub.space_before = Pt(10)

    links = [
        "• Live Web Application: http://localhost:5173",
        "• GitHub Repository: https://github.com/Sagheer1122/Arbisoft-GDG-Hackathon.git",
        "• Complete System PDF Manual: NurseFlow_Complete_Documentation.pdf",
        "• Kaggle Dataset Documentation: KAGGLE_DOCUMENTATION.md",
        "• REST API Documentation: API_DOCUMENTATION.md"
    ]
    for l in links:
        p_l = tf12.add_paragraph()
        p_l.text = l
        p_l.font.size = Pt(13)
        p_l.font.color.rgb = PURPLE_LIGHT
        p_l.space_before = Pt(10)

    # Save presentation
    output_path = "c:\\Users\\Tech Planet\\Desktop\\New folder (3)\\NurseFlow_Project_Presentation.pptx"
    prs.save(output_path)
    print(f"Professional presentation successfully saved at {output_path}")

if __name__ == "__main__":
    create_professional_presentation()
